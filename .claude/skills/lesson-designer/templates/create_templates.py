#!/usr/bin/env python3
"""
Create base templates for lesson designer skill.

Generates:
- slide_deck.pptx: PowerPoint template with branded layouts
- student_worksheet.docx: Word template with Jinja2 placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches
from docx.enum.text import WD_LINE_SPACING
import os

# Get script directory
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

def create_powerpoint_template():
    """Create PowerPoint template with proper layouts and formatting."""

    prs = Presentation()

    # The default presentation already has slide layouts
    # We'll verify and set default formatting

    # Add a title slide to establish formatting
    title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # Set title formatting (36pt minimum, per requirements)
    title.text = "Lesson Title"
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(44)
            run.font.bold = True

    # Set subtitle formatting
    subtitle.text = "Subtitle or Teacher Name"
    for paragraph in subtitle.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(24)

    # Add a content slide to establish body text formatting
    content_slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    # Set title formatting
    title.text = "Content Slide Title"
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(36)
            run.font.bold = True

    # Set body text formatting (20pt, above 16pt minimum per SLID-04)
    body.text = "Sample bullet point content\nSecond bullet point"
    for paragraph in body.text_frame.paragraphs:
        paragraph.level = 0
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(20)

    # Add a section header slide
    section_slide_layout = prs.slide_layouts[2]  # Section Header layout
    slide = prs.slides.add_slide(section_slide_layout)
    title = slide.shapes.title
    title.text = "Activity Transition"
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(40)
            run.font.bold = True

    # Add a blank slide for lesson plan
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Save template
    output_path = os.path.join(SCRIPT_DIR, 'slide_deck.pptx')
    prs.save(output_path)
    print(f"Created PowerPoint template: {output_path}")

    # Verify layouts
    print(f"  - Slide layouts available: {len(prs.slide_layouts)}")
    print(f"  - Sample slides created: {len(prs.slides)}")

    return output_path


def create_word_template():
    """Create Word template with Jinja2 placeholders."""

    doc = Document()

    # Set default font
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = DocxPt(12)

    # Set line spacing to 1.5 (will be enhanced in Phase 2)
    paragraph_format = style.paragraph_format
    paragraph_format.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
    paragraph_format.line_spacing = 1.5

    # Title
    title = doc.add_heading('STUDENT WORKSHEET', level=0)
    title.alignment = 1  # Center

    # Lesson information
    doc.add_paragraph()
    lesson_info = doc.add_paragraph()
    lesson_info.add_run('Lesson: ').bold = True
    lesson_info.add_run('{{ title }}')

    grade_info = doc.add_paragraph()
    grade_info.add_run('Grade ').bold = True
    grade_info.add_run('{{ grade_level }}')

    doc.add_paragraph()

    # Student info line
    student_info = doc.add_paragraph()
    student_info.add_run('Name: _________________    Date: _________________')

    doc.add_paragraph()

    # Learning objectives
    doc.add_heading('Learning Objectives:', level=2)
    objectives_template = """{% for objective in objectives %}
  {{ loop.index }}. {{ objective }}
{% endfor %}"""
    doc.add_paragraph(objectives_template)

    doc.add_paragraph()

    # Activities section
    activities_template = """{% for activity in activities %}
Activity {{ loop.index }}: {{ activity.name }}
Time: {{ activity.duration }} minutes

Instructions:
{% for instruction in activity.instructions %}
  {{ loop.index }}. {{ instruction }}
{% endfor %}

{% if activity.questions %}
Questions:
{% for question in activity.questions %}

{{ loop.index }}. {{ question }}

Answer: _________________________________________________

_________________________________________________________

_________________________________________________________

{% endfor %}
{% endif %}

{% endfor %}"""

    doc.add_paragraph(activities_template)

    doc.add_paragraph()

    # Vocabulary section
    vocab_template = """{% if vocabulary %}
Vocabulary:
{% for term in vocabulary %}
  {{ term.word }}: {{ term.definition }}
{% endfor %}
{% endif %}"""

    doc.add_paragraph(vocab_template)

    # Save template
    output_path = os.path.join(SCRIPT_DIR, 'student_worksheet.docx')
    doc.save(output_path)
    print(f"Created Word template: {output_path}")

    # Verify paragraphs
    print(f"  - Paragraphs created: {len(doc.paragraphs)}")
    print(f"  - Contains Jinja2 placeholders: {'{{' in doc.paragraphs[3].text}")

    return output_path


def verify_templates():
    """Verify that templates can be loaded by python-pptx and python-docx."""

    print("\nVerifying templates...")

    # Verify PowerPoint
    try:
        pptx_path = os.path.join(SCRIPT_DIR, 'slide_deck.pptx')
        prs = Presentation(pptx_path)
        assert len(prs.slide_layouts) >= 4, "Not enough slide layouts"
        print(f"✓ PowerPoint template valid ({len(prs.slide_layouts)} layouts)")
    except Exception as e:
        print(f"✗ PowerPoint template error: {e}")
        return False

    # Verify Word (python-docx can open it, docxtpl will use it)
    try:
        docx_path = os.path.join(SCRIPT_DIR, 'student_worksheet.docx')
        doc = Document(docx_path)
        has_placeholders = any('{{' in p.text for p in doc.paragraphs)
        assert has_placeholders, "No Jinja2 placeholders found"
        print(f"✓ Word template valid (has Jinja2 placeholders)")
    except Exception as e:
        print(f"✗ Word template error: {e}")
        return False

    print("\nAll templates created and verified successfully!")
    return True


if __name__ == '__main__':
    print("Creating lesson designer templates...\n")

    # Create templates
    create_powerpoint_template()
    create_word_template()

    # Verify
    verify_templates()
