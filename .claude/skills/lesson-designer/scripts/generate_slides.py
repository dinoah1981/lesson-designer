#!/usr/bin/env python3
"""
Generate Professional Slide Deck with Design System

Creates visually engaging, teacher-ready presentations using a professional
design system with coordinated colors, white content boxes, and activity icons.

Design principles based on Claude's natural presentation design capabilities.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import json
import sys
import os

# Design colors
HEADER_BLUE = RGBColor(0x2D, 0x5A, 0x87)
BODY_CHARCOAL = RGBColor(0x2C, 0x3E, 0x50)
ACCENT_GOLD = RGBColor(0xF4, 0xD0, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

# Activity icons by Marzano level
ICONS = {
    'retrieval': '🔄',
    'comprehension': '📚',
    'analysis': '🔍',
    'knowledge_utilization': '💡'
}

# Default discussion timing structure (in minutes)
DISCUSSION_STRUCTURE = {
    'opening': 2,      # Hook question, activate prior knowledge
    'pair_work': 5,    # Partner discussion
    'group_share': 5,  # Whole class sharing
    'closing': 3       # Synthesis, exit question
}

# Facilitation guidance templates
FACILITATION_PROMPTS = [
    "Can you build on what [student] said?",
    "Who has a different perspective?",
    "What evidence supports that claim?",
    "How does this connect to what we learned earlier?",
    "Can you give a specific example?"
]

DISCUSSION_WATCH_FOR = [
    "Students dominating airtime - redirect with 'Let's hear from someone new'",
    "Off-task conversations - physical proximity or refocus question",
    "Quiet students - invite with 'What do you think, [name]?'",
    "Surface-level responses - push with 'Tell me more about that'"
]


def is_discussion_activity(activity: dict) -> bool:
    """Detect if activity is a discussion based on name or type."""
    name_lower = activity.get('name', '').lower()
    discussion_keywords = ['discussion', 'debate', 'seminar', 'share',
                           'pair', 'group talk', 'debrief', 'reflection']
    return any(kw in name_lower for kw in discussion_keywords)


def add_header_bar(prs, slide, title_text, subtitle_text=None, timer_text=None):
    """Add professional header bar with title."""
    # Header background
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = HEADER_BLUE
    header.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = title_text
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = WHITE
    title_run.font.italic = True

    # Subtitle if provided
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(10), Inches(0.4))
        sub_frame = sub_box.text_frame
        sub_p = sub_frame.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle_text
        sub_run.font.size = Pt(20)
        sub_run.font.color.rgb = WHITE

    # Timer if provided
    if timer_text:
        timer_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.5), Inches(1.5), Inches(0.5))
        timer_frame = timer_box.text_frame
        timer_p = timer_frame.paragraphs[0]
        timer_p.alignment = PP_ALIGN.RIGHT
        timer_run = timer_p.add_run()
        timer_run.text = timer_text
        timer_run.font.size = Pt(24)
        timer_run.font.bold = True
        timer_run.font.color.rgb = ACCENT_GOLD


def add_content_box(prs, slide, left, top, width, height):
    """Add white content rectangle with subtle border."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    box.line.width = Pt(1)
    # Adjust corner radius
    box.adjustments[0] = 0.05
    return box


def add_light_background(prs, slide):
    """Add light gray background below header."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.5),
        prs.slide_width, Inches(6)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()
    return bg


def create_hidden_lesson_plan(prs, slide, lesson):
    """Create hidden first slide with lesson plan for teacher."""
    slide._element.set('show', '0')  # HIDE this slide

    # Get slide content - check both new schema (slide_content) and legacy (hidden_slide_content)
    slide_content = lesson.get('slide_content', lesson.get('hidden_slide_content', {}))

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '📋 LESSON PLAN (Teacher Reference Only)'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = HEADER_BLUE

    # Objective
    obj_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12), Inches(0.8))
    tf = obj_box.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = 'OBJECTIVE: '
    run1.font.bold = True
    run1.font.size = Pt(16)
    run2 = p.add_run()
    run2.text = lesson['objective']
    run2.font.size = Pt(16)

    # Agenda
    agenda_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(2.5))
    tf = agenda_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'AGENDA:'
    run.font.bold = True
    run.font.size = Pt(16)
    for item in slide_content.get('agenda', []):
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  • {item['activity']} ({item['duration']} min)"
        run.font.size = Pt(14)

    # Misconceptions
    misc_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.8), Inches(6), Inches(2.5))
    tf = misc_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'WATCH FOR (Misconceptions):'
    run.font.bold = True
    run.font.size = Pt(16)
    for misc in slide_content.get('misconceptions', [])[:3]:
        p = tf.add_paragraph()
        run = p.add_run()
        text = f"  ⚠️ {misc[:80]}..." if len(misc) > 80 else f"  ⚠️ {misc}"
        run.text = text
        run.font.size = Pt(12)

    # Delivery tips
    tips_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2.5))
    tf = tips_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'DELIVERY TIPS:'
    run.font.bold = True
    run.font.size = Pt(16)
    for tip in slide_content.get('delivery_tips', [])[:4]:
        p = tf.add_paragraph()
        run = p.add_run()
        text = f"  💡 {tip[:100]}" if len(tip) > 100 else f"  💡 {tip}"
        run.text = text
        run.font.size = Pt(12)


def create_title_slide(prs, slide, lesson):
    """Create professional title slide."""
    # Full blue background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = HEADER_BLUE
    bg.line.fill.background()

    # White content area
    add_content_box(prs, slide, Inches(1), Inches(1.5), Inches(11.333), Inches(4.5))

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = lesson['title']
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = BODY_CHARCOAL

    # Subtitle info
    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.333), Inches(0.6))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{lesson['grade_level']}  •  {lesson['duration']} minutes"
    run.font.size = Pt(24)
    run.font.color.rgb = BODY_CHARCOAL

    # Objective
    obj_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.333), Inches(1))
    tf = obj_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    obj_text = lesson['objective'][:100] if len(lesson['objective']) > 100 else lesson['objective']
    run.text = f"🎯 {obj_text}"
    run.font.size = Pt(18)
    run.font.color.rgb = BODY_CHARCOAL

    # Gold accent bar at bottom
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8),
        prs.slide_width, Inches(0.7)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_GOLD
    accent.line.fill.background()


def create_agenda_slide(prs, slide, lesson):
    """Create agenda slide with activity checklist."""
    add_header_bar(prs, slide, "📋 Today's Agenda")
    add_light_background(prs, slide)

    # Agenda items with checkboxes
    y_pos = 2.0
    for activity in lesson['activities']:
        icon = ICONS.get(activity['marzano_level'], '📌')

        # Item box
        add_content_box(prs, slide, Inches(1), Inches(y_pos), Inches(11.333), Inches(0.9))

        # Activity name
        text_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos + 0.15), Inches(9), Inches(0.6))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"☐ {icon} {activity['name']}"
        run.font.size = Pt(22)
        run.font.color.rgb = BODY_CHARCOAL

        # Duration
        dur_box = slide.shapes.add_textbox(Inches(10.5), Inches(y_pos + 0.15), Inches(1.5), Inches(0.6))
        tf = dur_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{activity['duration']} min"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = ACCENT_GOLD

        y_pos += 1.0


def create_discussion_slide(prs, slide, activity):
    """Create structured discussion slide with timing breakdown."""
    duration = activity.get('duration', 15)

    # Header with total time
    add_header_bar(prs, slide,
        f"💬 {activity['name']}",
        subtitle_text="Discussion Activity",
        timer_text=f"⏱️ {duration}m")

    add_light_background(prs, slide)

    # Calculate timing breakdown
    opening_time = min(2, duration // 6)
    closing_time = min(3, duration // 5)
    main_time = duration - opening_time - closing_time
    pair_time = main_time // 2
    share_time = main_time - pair_time

    # Content sections with timing
    y_pos = 1.9

    # Opening section
    add_content_box(prs, slide, Inches(0.5), Inches(y_pos), Inches(12.333), Inches(1.2))
    opening_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(11.9), Inches(1))
    tf = opening_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"🎯 Opening ({opening_time} min)"
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = HEADER_BLUE

    # Get opening - prefer student_directions over teacher instructions
    student_directions = activity.get('student_directions', '')
    opening = activity.get('discussion_structure', {}).get('opening', student_directions or 'Think about the topic...')
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = opening[:100] if len(str(opening)) > 100 else str(opening)
    run2.font.size = Pt(16)

    y_pos += 1.4

    # Main discussion section
    add_content_box(prs, slide, Inches(0.5), Inches(y_pos), Inches(12.333), Inches(2.5))
    main_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(11.9), Inches(2.3))
    tf = main_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"💭 Discussion ({main_time} min: Pair {pair_time}m | Share {share_time}m)"
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = HEADER_BLUE

    # Discussion prompts - prefer student_questions or discussion_questions
    prompts = activity.get('discussion_structure', {}).get('prompts',
        activity.get('student_questions', activity.get('discussion_questions', ['What do you think about this topic?'])))
    for prompt in prompts[:3]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {prompt[:80]}"
        run.font.size = Pt(16)

    y_pos += 2.7

    # Closing section
    add_content_box(prs, slide, Inches(0.5), Inches(y_pos), Inches(12.333), Inches(1.1))
    closing_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(11.9), Inches(0.9))
    tf = closing_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"✅ Closing ({closing_time} min)"
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = HEADER_BLUE

    closing = activity.get('discussion_structure', {}).get('closing', 'Synthesize key takeaways')
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = closing[:100] if len(str(closing)) > 100 else str(closing)
    run2.font.size = Pt(16)

    # Full facilitation notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = f"""FACILITATION GUIDE: {activity['name']}

TIME ALLOCATION ({duration} min total):
  • Opening: {opening_time} min
  • Pair discussion: {pair_time} min
  • Whole group share: {share_time} min
  • Closing: {closing_time} min

FACILITATION MOVES:
  • During pair work: Circulate, listen, note interesting ideas to highlight
  • During share: Track who speaks, balance airtime across students
  • Productive struggle: Let silence sit for 3-5 seconds before intervening

PROMPTS TO USE:
"""
    for prompt in FACILITATION_PROMPTS[:4]:
        notes_tf.text += f"  • '{prompt}'\n"

    notes_tf.text += "\nWATCH FOR:\n"
    for warning in DISCUSSION_WATCH_FOR[:3]:
        notes_tf.text += f"  ⚠️ {warning}\n"

    # Add activity-specific notes if provided
    if activity.get('assessment_method'):
        notes_tf.text += f"\nASSESSMENT: {activity['assessment_method']}\n"


def create_activity_slide(prs, slide, activity):
    """Create slide for individual activity with student-facing content."""
    icon = ICONS.get(activity.get('marzano_level', 'comprehension'), '📌')

    add_header_bar(
        prs, slide,
        f"{icon} {activity['name']}",
        subtitle_text=activity.get('marzano_level', 'comprehension').replace('_', ' ').title(),
        timer_text=f"⏱️ {activity.get('duration', 10)}m"
    )

    add_light_background(prs, slide)

    # Get student-facing content
    student_directions = activity.get('student_directions', '')
    student_questions = activity.get('student_questions', [])
    discussion_questions = activity.get('discussion_questions', [])
    visual_description = activity.get('visual_description', '')
    recommended_image = activity.get('recommended_image', {})

    # Check if we have a visual (either description or recommended image)
    has_visual = visual_description or recommended_image
    img_desc = recommended_image.get('description', visual_description) or visual_description
    img_url = recommended_image.get('url', '')

    # Layout depends on whether we have a visual
    if has_visual:
        # Split layout: image placeholder on left, content on right
        # Image placeholder box
        add_content_box(prs, slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(4.5))
        img_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5), Inches(2.5))
        tf = img_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"[IMAGE: {img_desc[:60]}]"
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = RGBColor(128, 128, 128)

        # Add URL if available
        if img_url:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = img_url[:50] + "..." if len(img_url) > 50 else img_url
            run2.font.size = Pt(10)
            run2.font.color.rgb = RGBColor(0, 102, 204)  # Blue link color

        # Content on right side
        add_content_box(prs, slide, Inches(6.2), Inches(1.8), Inches(6.633), Inches(4.5))
        content_left = 6.5
        content_width = 6
    else:
        # Full width content
        add_content_box(prs, slide, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))
        content_left = 0.8
        content_width = 11.7

    # Content text box
    content_box = slide.shapes.add_textbox(Inches(content_left), Inches(2.1), Inches(content_width), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Add student directions first
    if student_directions:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = student_directions[:150] if len(student_directions) > 150 else student_directions
        run.font.size = Pt(20)
        run.font.color.rgb = BODY_CHARCOAL
        p.space_after = Pt(18)

    # Add student questions or discussion questions
    questions = student_questions or discussion_questions
    if questions:
        for i, q in enumerate(questions[:4]):  # Max 4 questions per slide
            p = tf.add_paragraph() if (student_directions or i > 0) else tf.paragraphs[0]
            run = p.add_run()
            q_text = q[:80] if len(q) > 80 else q
            run.text = f"{i+1}. {q_text}"
            run.font.size = Pt(18)
            run.font.color.rgb = BODY_CHARCOAL
            p.space_after = Pt(12)

    # If no student content, fall back to a generic prompt
    if not student_directions and not questions:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Complete the {activity['name']} activity."
        run.font.size = Pt(24)
        run.font.color.rgb = BODY_CHARCOAL

    # Add presenter notes with TEACHER instructions
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = f"TEACHER GUIDE: {activity['name']}\n\n"
    notes_tf.text += f"TIME: {activity.get('duration', 10)} minutes\n\n"

    # Add image URL if available
    if img_url:
        notes_tf.text += f"IMAGE URL: {img_url}\n\n"

    # Teacher moves go in notes
    teacher_moves = activity.get('teacher_moves', [])
    if teacher_moves:
        notes_tf.text += "TEACHER MOVES:\n"
        for move in teacher_moves:
            notes_tf.text += f"- {move}\n"
        notes_tf.text += "\n"

    # Instructions (teacher steps) go in notes
    instructions = activity.get('instructions', [])
    if instructions:
        notes_tf.text += "STEP-BY-STEP:\n"
        for i, instr in enumerate(instructions):
            notes_tf.text += f"{i+1}. {instr}\n"
        notes_tf.text += "\n"

    notes_tf.text += f"ASSESSMENT: {activity.get('assessment_method', 'Observe student work')}\n"

    # Differentiation notes
    diff = activity.get('differentiation', {})
    if diff:
        notes_tf.text += f"\nSUPPORT: {diff.get('support', 'Provide scaffolding')}\n"
        notes_tf.text += f"EXTENSION: {diff.get('extension', 'Challenge with deeper questions')}\n"


def create_vocabulary_slide(prs, slide, vocab):
    """Create vocabulary slide with terms and definitions."""
    add_header_bar(prs, slide, "📖 Key Vocabulary")
    add_light_background(prs, slide)

    y_pos = 1.8
    for term in vocab[:5]:
        # Term box
        add_content_box(prs, slide, Inches(0.5), Inches(y_pos), Inches(3), Inches(0.9))

        # Term text
        t_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.2), Inches(2.6), Inches(0.5))
        tf = t_box.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = term['word']
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = HEADER_BLUE

        # Definition box
        add_content_box(prs, slide, Inches(3.7), Inches(y_pos), Inches(9.133), Inches(0.9))

        # Definition text
        d_box = slide.shapes.add_textbox(Inches(3.9), Inches(y_pos + 0.15), Inches(8.7), Inches(0.6))
        tf = d_box.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        def_text = term['definition'][:80] if len(term['definition']) > 80 else term['definition']
        run.text = def_text
        run.font.size = Pt(16)
        run.font.color.rgb = BODY_CHARCOAL

        y_pos += 1.0


def create_key_facts_slide(prs, slide, key_facts):
    """Create a slide showing key facts students should know."""
    add_header_bar(prs, slide, "📋 Key Facts")
    add_light_background(prs, slide)

    add_content_box(prs, slide, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))

    facts_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.5))
    tf = facts_box.text_frame
    tf.word_wrap = True

    for i, fact in enumerate(key_facts[:8]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {fact}"
        run.font.size = Pt(18)
        run.font.color.rgb = BODY_CHARCOAL
        p.space_after = Pt(8)


def create_content_table_slide(prs, slide, content_table, title="Reference Data"):
    """Create a slide with a data table."""
    add_header_bar(prs, slide, f"📊 {title}")
    add_light_background(prs, slide)

    headers = content_table.get('headers', [])
    rows = content_table.get('rows', [])

    if not headers or not rows:
        return

    # Calculate table dimensions
    num_cols = len(headers)
    num_rows = min(len(rows) + 1, 8)  # Header + data rows, max 8 total

    # Create table shape
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5)).table

    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BLUE
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE

    # Fill data rows
    for row_idx, row_data in enumerate(rows[:num_rows-1]):
        for col_idx, cell_value in enumerate(row_data[:num_cols]):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_value)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = BODY_CHARCOAL


def create_exit_ticket_slide(prs, slide, assessment):
    """Create exit ticket slide with multiple choice and short answer questions."""
    add_header_bar(prs, slide, "🎫 Exit Ticket")
    add_light_background(prs, slide)

    add_content_box(prs, slide, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))

    q_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.5))
    tf = q_box.text_frame
    tf.word_wrap = True

    q_num = 1
    first_para = True

    # Multiple choice questions
    mc_questions = assessment.get('multiple_choice', [])
    for mc in mc_questions[:2]:  # Max 2 MC on one slide
        question = mc.get('question', '')
        options = mc.get('options', [])

        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.text = f"{q_num}. {question}"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = BODY_CHARCOAL
        p.space_after = Pt(4)

        # Options
        option_letters = ['a', 'b', 'c', 'd']
        for i, option in enumerate(options[:4]):
            p = tf.add_paragraph()
            run = p.add_run()
            letter = option_letters[i] if i < len(option_letters) else str(i+1)
            run.text = f"    {letter}. {option}"
            run.font.size = Pt(14)
            run.font.color.rgb = BODY_CHARCOAL
            p.space_after = Pt(2)

        q_num += 1

    # Short answer questions
    short_answer = assessment.get('short_answer', [])
    for sa in short_answer[:2]:  # Max 2 SA on one slide
        question = sa.get('question', '')

        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.text = f"{q_num}. {question}"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = BODY_CHARCOAL
        p.space_after = Pt(12)
        q_num += 1

    # Fallback to old format
    questions = assessment.get('questions', [])
    if not mc_questions and not short_answer and questions:
        for i, q in enumerate(questions[:3]):
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"{i+1}. {q}"
            run.font.size = Pt(20)
            run.font.color.rgb = BODY_CHARCOAL
            p.space_after = Pt(18)


def generate_slides(lesson_path: str, output_path: str) -> bool:
    """Generate professional slide deck from lesson design."""

    # Load lesson data
    with open(lesson_path, 'r', encoding='utf-8') as f:
        lesson = json.load(f)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Hidden lesson plan
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    create_hidden_lesson_plan(prs, slide1, lesson)

    # Slide 2: Title slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    create_title_slide(prs, slide2, lesson)

    # Slide 3: Agenda
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    create_agenda_slide(prs, slide3, lesson)

    # Key facts slide (if key_facts exists) - show early so students have context
    key_facts = lesson.get('key_facts', [])
    if key_facts:
        slide_facts = prs.slides.add_slide(prs.slide_layouts[6])
        create_key_facts_slide(prs, slide_facts, key_facts)

    # Content table slide (if content_table exists) - reference data for students
    content_table = lesson.get('content_table', {})
    if content_table.get('headers') and content_table.get('rows'):
        slide_table = prs.slides.add_slide(prs.slide_layouts[6])
        create_content_table_slide(prs, slide_table, content_table)

    # Vocabulary slide (if vocabulary exists)
    vocab = lesson.get('vocabulary', [])
    if vocab:
        slide_vocab = prs.slides.add_slide(prs.slide_layouts[6])
        create_vocabulary_slide(prs, slide_vocab, vocab)

    # Activity slides
    for activity in lesson['activities']:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if is_discussion_activity(activity):
            create_discussion_slide(prs, slide, activity)
        else:
            create_activity_slide(prs, slide, activity)

    # Exit ticket slide
    assessment = lesson.get('assessment', {})
    if assessment.get('questions') or assessment.get('multiple_choice') or assessment.get('short_answer'):
        slide_exit = prs.slides.add_slide(prs.slide_layouts[6])
        create_exit_ticket_slide(prs, slide_exit, assessment)

    # Save presentation
    prs.save(output_path)
    print(f"Created professional slide deck: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return True


def main():
    """CLI entry point."""
    if len(sys.argv) < 2:
        print("Usage: python create_professional_slides.py <lesson.json> [output.pptx]")
        sys.exit(1)

    lesson_path = sys.argv[1]

    if not os.path.exists(lesson_path):
        print(f"Error: Lesson file not found: {lesson_path}", file=sys.stderr)
        sys.exit(1)

    # Determine output path
    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
    else:
        lesson_dir = os.path.dirname(os.path.abspath(lesson_path))
        output_path = os.path.join(lesson_dir, '05_slides.pptx')

    success = generate_slides(lesson_path, output_path)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
