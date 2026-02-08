"""
Lesson Designer v2 — Document Templates
Deterministic template functions for professional slide decks and worksheets.
Matches the visual design system from Claude.ai skill output.
"""

import io
import re
import tempfile
from typing import List, Dict, Tuple, Optional

import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, Twips, RGBColor as DocxRGBColor, Emu as DocxEmu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn as docx_qn
from docx.oxml import OxmlElement


# ═══════════════════════════════════════════════════════════════════════════════
# COLOR PALETTE
# ═══════════════════════════════════════════════════════════════════════════════

# Slide colors
S_DARK_BG = "1A1F36"
S_NAVY = "2D3561"
S_TEAL = "1B998B"
S_GOLD = "C69749"
S_TERRA = "B85042"
S_CREAM = "F7F7F2"
S_LIGHT = "EEF0F6"
S_WHITE = "FFFFFF"
S_TEXT = "1D1D1D"
S_GRAY = "6B7280"

# Doc colors
D_NAVY = "2D3561"
D_WHITE = "FFFFFF"
D_LAVENDER = "EEF0F6"
D_TEAL = "1B998B"
D_TEXT = "1D1D1D"
D_GRAY = "6B7280"
D_BORDER = "CCCCCC"

# Fonts
SERIF = "Georgia"
SANS = "Helvetica"

# Slide dimensions (10" x 5.625" widescreen 16:9)
SW = Inches(10)
SH = Inches(5.625)


# ═══════════════════════════════════════════════════════════════════════════════
# CHART RENDERING (matplotlib → image bytes)
# ═══════════════════════════════════════════════════════════════════════════════

# Chart color palette (matches slide design system)
_CHART_COLORS = ["#1B998B", "#2D3561", "#C69749", "#B85042", "#4DAE58", "#6B7280"]
_CHART_BG = "#F7F7F2"
_CHART_TEXT = "#1D1D1D"
_CHART_GRID = "#CCCCCC"


def _parse_chart_spec(block: str) -> Optional[Dict]:
    """Parse a ```chart``` code block into a spec dict."""
    spec = {}
    for line in block.strip().split("\n"):
        line = line.strip()
        if ":" not in line:
            continue
        key, _, val = line.partition(":")
        key = key.strip().lower()
        val = val.strip()
        # Parse list values: [a, b, c]
        if val.startswith("[") and val.endswith("]"):
            items = [v.strip().strip("'\"") for v in val[1:-1].split(",")]
            # Try to convert to numbers
            parsed = []
            for item in items:
                try:
                    parsed.append(float(item))
                except ValueError:
                    parsed.append(item)
            spec[key] = parsed
        else:
            try:
                spec[key] = float(val)
            except ValueError:
                spec[key] = val
    return spec if spec.get("type") else None


def _render_chart(spec: Dict) -> Optional[bytes]:
    """Render a chart spec to PNG bytes using matplotlib."""
    chart_type = spec.get("type", "").lower()
    title = spec.get("title", "")

    fig, ax = plt.subplots(figsize=(8, 4.5), dpi=150)
    fig.patch.set_facecolor(_CHART_BG)
    ax.set_facecolor(_CHART_BG)
    ax.tick_params(colors=_CHART_TEXT, labelsize=10)
    ax.set_title(title, fontsize=14, fontweight='bold', color=_CHART_TEXT, pad=12)

    try:
        if chart_type in ("line", "scatter"):
            x_vals = spec.get("x_values", [])
            if not x_vals:
                plt.close(fig)
                return None
            x = np.array([float(v) for v in x_vals])

            # Render multiple series
            series_idx = 0
            for suffix in ["", "2", "3", "4"]:
                y_key = f"y_values{suffix}" if suffix else "y_values"
                label_key = f"series_label{suffix}" if suffix else "series_label"
                if y_key not in spec:
                    continue
                y = np.array([float(v) for v in spec[y_key]])
                label = spec.get(label_key, f"Series {series_idx + 1}")
                color = _CHART_COLORS[series_idx % len(_CHART_COLORS)]

                if chart_type == "line":
                    # Plot smooth curve if enough points
                    if len(x) >= 4:
                        x_smooth = np.linspace(x.min(), x.max(), 200)
                        coeffs = np.polyfit(x, y, min(len(x) - 1, 6))
                        y_smooth = np.polyval(coeffs, x_smooth)
                        ax.plot(x_smooth, y_smooth, color=color, linewidth=2.5, label=label)
                    else:
                        ax.plot(x, y, color=color, linewidth=2.5, label=label)
                    ax.scatter(x, y, color=color, s=40, zorder=5)
                else:
                    ax.scatter(x, y, color=color, s=60, label=label, zorder=5)
                series_idx += 1

            ax.axhline(y=0, color=_CHART_GRID, linewidth=0.8, zorder=1)
            ax.axvline(x=0, color=_CHART_GRID, linewidth=0.8, zorder=1)
            ax.grid(True, alpha=0.3, color=_CHART_GRID)
            if series_idx > 1:
                ax.legend(fontsize=10, loc='best')

            # Annotate labeled points
            labels = spec.get("labels", spec.get("point_labels", []))
            if isinstance(labels, list):
                pts_x = spec.get("x_values", [])
                pts_y = spec.get("y_values", [])
                for i, lbl in enumerate(labels):
                    if i < len(pts_x) and i < len(pts_y) and str(lbl).strip():
                        ax.annotate(str(lbl), (float(pts_x[i]), float(pts_y[i])),
                                    textcoords="offset points", xytext=(0, 12),
                                    ha='center', fontsize=9, color=_CHART_TEXT,
                                    fontweight='bold',
                                    bbox=dict(boxstyle='round,pad=0.3',
                                              facecolor='white', edgecolor=_CHART_GRID))

        elif chart_type == "bar":
            categories = [str(c) for c in spec.get("categories", [])]
            values = [float(v) for v in spec.get("values", [])]
            if not categories or not values:
                plt.close(fig)
                return None
            colors = [_CHART_COLORS[i % len(_CHART_COLORS)] for i in range(len(categories))]
            bars = ax.bar(categories, values, color=colors, width=0.6, edgecolor='white')
            for bar, val in zip(bars, values):
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(values)*0.02,
                        f'{val:g}', ha='center', va='bottom', fontsize=10,
                        fontweight='bold', color=_CHART_TEXT)
            ax.set_ylabel(spec.get("series_label", ""), fontsize=11, color=_CHART_TEXT)
            ax.grid(axis='y', alpha=0.3, color=_CHART_GRID)

        elif chart_type == "number_line":
            mn = float(spec.get("min", -10))
            mx = float(spec.get("max", 10))
            points = [float(p) for p in spec.get("points", [])]
            point_labels = [str(l) for l in spec.get("point_labels", [])]

            fig, ax = plt.subplots(figsize=(8, 2), dpi=150)
            fig.patch.set_facecolor(_CHART_BG)
            ax.set_facecolor(_CHART_BG)
            ax.set_title(title, fontsize=14, fontweight='bold', color=_CHART_TEXT, pad=8)

            ax.arrow(mn - 0.5, 0, (mx - mn + 1), 0, head_width=0.15, head_length=0.2,
                     fc=_CHART_TEXT, ec=_CHART_TEXT, linewidth=1.5)
            ax.set_xlim(mn - 1, mx + 1)
            ax.set_ylim(-0.8, 0.8)

            # Tick marks
            for t in range(int(mn), int(mx) + 1):
                ax.plot([t, t], [-0.1, 0.1], color=_CHART_TEXT, linewidth=1)
                ax.text(t, -0.3, str(t), ha='center', fontsize=9, color=_CHART_TEXT)

            # Highlighted points
            for i, p in enumerate(points):
                ax.plot(p, 0, 'o', markersize=12, color=_CHART_COLORS[0], zorder=5)
                lbl = point_labels[i] if i < len(point_labels) else ""
                if lbl:
                    ax.text(p, 0.4, lbl, ha='center', fontsize=10, fontweight='bold',
                            color=_CHART_COLORS[0])

            ax.axis('off')

        elif chart_type == "comparison":
            left_title = spec.get("left_title", "A")
            right_title = spec.get("right_title", "B")
            left_items = [str(i) for i in spec.get("left_items", [])]
            right_items = [str(i) for i in spec.get("right_items", [])]

            fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(8, 4), dpi=150)
            fig.patch.set_facecolor(_CHART_BG)
            fig.suptitle(title, fontsize=14, fontweight='bold', color=_CHART_TEXT, y=0.98)

            for ax_panel, panel_title, items, color in [
                (ax1, left_title, left_items, _CHART_COLORS[0]),
                (ax2, right_title, right_items, _CHART_COLORS[1]),
            ]:
                ax_panel.set_facecolor('white')
                ax_panel.set_xlim(0, 1)
                ax_panel.set_ylim(0, 1)
                ax_panel.axis('off')

                # Title bar
                ax_panel.add_patch(plt.Rectangle((0, 0.85), 1, 0.15,
                                                  facecolor=color, transform=ax_panel.transAxes))
                ax_panel.text(0.5, 0.92, panel_title, ha='center', va='center',
                              fontsize=13, fontweight='bold', color='white',
                              transform=ax_panel.transAxes)

                # Items
                n = len(items)
                for j, item in enumerate(items):
                    y_pos = 0.78 - j * (0.75 / max(n, 1))
                    ax_panel.text(0.08, y_pos, f"\u2022 {item}", fontsize=10,
                                  va='top', color=_CHART_TEXT, transform=ax_panel.transAxes,
                                  wrap=True)

                # Border
                for spine in ax_panel.spines.values():
                    spine.set_visible(True)
                    spine.set_color(_CHART_GRID)

            plt.tight_layout(rect=[0, 0, 1, 0.95])

        else:
            plt.close(fig)
            return None

        if chart_type not in ("comparison",):
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color(_CHART_GRID)
            ax.spines['bottom'].set_color(_CHART_GRID)

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight', facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()

    except Exception:
        plt.close(fig)
        return None


def _extract_chart_blocks(text: str) -> Tuple[str, List[Dict]]:
    """Extract ```chart``` blocks from content text. Returns (cleaned_text, chart_specs)."""
    charts = []
    cleaned = text
    for match in re.finditer(r'```chart\s*\n(.*?)```', text, re.DOTALL):
        spec = _parse_chart_spec(match.group(1))
        if spec:
            charts.append(spec)
        cleaned = cleaned.replace(match.group(0), "")
    return cleaned.strip(), charts


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _drgb(h):
    return DocxRGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.fill.background()
    return shape


def _txt(slide, text, x, y, w, h, font=SANS, sz=14, bold=False, italic=False,
         color=S_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.name = font
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = _rgb(color)
    p.alignment = align
    return tb


def _add_runs(p, text, font=SANS, sz=14, color=S_TEXT, base_bold=False):
    """Add text with **bold** markdown support to a paragraph."""
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if not part:
            continue
        run = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        else:
            run.text = part
            if base_bold:
                run.font.bold = True
        run.font.name = font
        run.font.size = Pt(sz)
        run.font.color.rgb = _rgb(color)


def _rich_txt(slide, text, x, y, w, h, font=SANS, sz=14, color=S_TEXT,
              align=PP_ALIGN.LEFT):
    """Like _txt but parses **bold** markdown into actual bold formatting."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _add_runs(p, str(text), font=font, sz=sz, color=color)
    return tb


def _bullets(slide, lines, x, y, w, h, font=SANS, sz=18, color=S_TEXT,
             spacing=1.2, bullet_char=True):
    """Add a text box with multiple bullet lines."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "\u2022 " if bullet_char else ""
        _add_runs(p, prefix + line, font=font, sz=sz, color=color)
        p.space_after = Pt(sz * (spacing - 1) + 3)
        p.alignment = PP_ALIGN.LEFT

    return tb


def _header_bar(slide, title, fill=S_NAVY, time_str=None):
    """Colored header bar across top of slide with title."""
    _rect(slide, Inches(0), Inches(0), SW, Inches(0.9), fill)
    _txt(slide, title, Inches(0.5), Inches(0.15), Inches(8.5), Inches(0.6),
         font=SERIF, sz=22, bold=True, color=S_WHITE)
    if time_str:
        _txt(slide, time_str, Inches(8.3), Inches(0.2), Inches(1.4), Inches(0.5),
             font=SANS, sz=16, italic=True, color=S_GOLD, align=PP_ALIGN.RIGHT)


def _notes(slide, text):
    """Set speaker notes."""
    if text:
        slide.notes_slide.notes_text_frame.text = text[:3000]


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE LAYOUT FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def _blank(prs):
    """Add a blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_title(prs, title, subtitle=""):
    """Dark background title slide."""
    s = _blank(prs)
    _bg(s, S_DARK_BG)
    _rect(s, Inches(0), Inches(0), SW, Inches(0.06), S_TEAL)
    _txt(s, title, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.2),
         font=SERIF, sz=40, bold=True, color=S_WHITE)
    if subtitle:
        _txt(s, subtitle, Inches(0.8), Inches(2.8), Inches(8.4), Inches(0.6),
             font=SANS, sz=18, color=S_GOLD)
    _rect(s, Inches(0.8), Inches(4.9), Inches(2), Inches(0.04), S_TEAL)
    return s


def slide_hidden_plan(prs, plan_text):
    """Hidden teacher reference slide with full lesson plan."""
    s = _blank(prs)
    _bg(s, S_DARK_BG)
    _txt(s, "LESSON PLAN \u2014 TEACHER REFERENCE",
         Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.4),
         font=SERIF, sz=14, bold=True, color=S_GOLD)
    _txt(s, plan_text[:3500], Inches(0.3), Inches(0.5), Inches(9.4), Inches(4.9),
         font=SANS, sz=8, color=S_WHITE)
    s._element.set('show', '0')
    return s


def slide_objectives(prs, objectives, duration):
    """Learning objectives slide with accent cards."""
    s = _blank(prs)
    _bg(s, S_CREAM)
    _header_bar(s, "Today\u2019s Learning Goals")

    n = len(objectives)
    card_h = min(0.85, (3.8 - 0.2 * n) / max(n, 1))
    y = 1.15

    for i, obj in enumerate(objectives):
        _rect(s, Inches(0.6), Inches(y), Inches(0.08), Inches(card_h), S_TEAL)
        _rect(s, Inches(0.68), Inches(y), Inches(8.72), Inches(card_h), S_WHITE)
        label = f"Objective {i+1}:" if n > 1 else "Objective:"
        _txt(s, label, Inches(0.9), Inches(y + 0.05), Inches(1.8), Inches(0.25),
             font=SANS, sz=12, bold=True, color=S_TEAL)
        _txt(s, obj, Inches(0.9), Inches(y + 0.28), Inches(8.3), Inches(card_h - 0.33),
             font=SANS, sz=16, color=S_TEXT)
        y += card_h + 0.2

    _txt(s, f"{duration} minutes", Inches(7.5), Inches(min(y + 0.1, 4.9)), Inches(2), Inches(0.4),
         font=SANS, sz=14, italic=True, color=S_GOLD, align=PP_ALIGN.RIGHT)
    return s


def slide_content(prs, title, items, notes_text="", time_str=None,
                  bar_color=S_NAVY, bg_color=S_CREAM):
    """Standard content slide: header bar + bullet points + optional chart."""
    s = _blank(prs)
    _bg(s, bg_color)
    _header_bar(s, title, fill=bar_color, time_str=time_str)
    if items:
        _bullets(s, items[:5], Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0),
                 font=SANS, sz=16, color=S_TEXT)
    _notes(s, notes_text)
    return s


def slide_visual(prs, chart_spec, bar_color=S_NAVY, bg_color=S_CREAM):
    """Render a chart spec as an image slide."""
    img_bytes = _render_chart(chart_spec)
    if not img_bytes:
        return None
    s = _blank(prs)
    _bg(s, bg_color)
    title = chart_spec.get("title", "")
    if title:
        _header_bar(s, title, fill=bar_color)
        img_top = Inches(1.0)
        img_h = Inches(4.4)
    else:
        img_top = Inches(0.2)
        img_h = Inches(5.2)
    # Center the image
    img_stream = io.BytesIO(img_bytes)
    s.shapes.add_picture(img_stream, Inches(0.6), img_top, Inches(8.8), img_h)
    return s


def slide_table(prs, title, headers, rows, bar_color=S_NAVY, bg_color=S_CREAM):
    """Content slide with a data table (for coordinate tables, comparisons, etc.)."""
    s = _blank(prs)
    _bg(s, bg_color)
    _header_bar(s, title, fill=bar_color)

    n_cols = len(headers)
    n_rows = min(len(rows), 8)
    table_w = min(n_cols * 1.5, 8.8)
    x_off = (10 - table_w) / 2  # center the table
    col_w = table_w / n_cols

    tbl = s.shapes.add_table(1 + n_rows, n_cols,
                              Inches(x_off), Inches(1.2),
                              Inches(table_w), Inches(0.4 + n_rows * 0.35)).table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.name = SANS
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = _rgb(S_WHITE)
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(S_NAVY)

    # Data rows
    for i, row_data in enumerate(rows[:n_rows]):
        fill = S_LIGHT if i % 2 == 0 else S_WHITE
        for j in range(n_cols):
            cell = tbl.cell(i + 1, j)
            cell.text = row_data[j] if j < len(row_data) else ""
            for p in cell.text_frame.paragraphs:
                p.font.name = SANS
                p.font.size = Pt(11)
                p.font.color.rgb = _rgb(S_TEXT)
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(fill)

    return s


def slide_cards(prs, title, cards, bar_color=S_NAVY, bg_color=S_CREAM):
    """Info cards with colored left accent bars."""
    s = _blank(prs)
    _bg(s, bg_color)
    _header_bar(s, title, fill=bar_color)

    accents = [S_TEAL, S_GOLD, S_TERRA, S_NAVY]
    n = len(cards[:4])
    card_h = min(1.0, max(0.65, 3.8 / max(n, 1)))
    y = 1.1

    for i, (ct, cb) in enumerate(cards[:4]):
        c = accents[i % len(accents)]
        _rect(s, Inches(0.6), Inches(y), Inches(0.08), Inches(card_h), c)
        _rect(s, Inches(0.68), Inches(y), Inches(8.72), Inches(card_h), S_WHITE)
        _txt(s, ct, Inches(0.9), Inches(y + 0.05), Inches(8.3), Inches(0.25),
             font=SANS, sz=14, bold=True, color=c)
        _txt(s, cb[:180], Inches(0.9), Inches(y + 0.30), Inches(8.3), Inches(card_h - 0.35),
             font=SANS, sz=12, color=S_TEXT)
        y += card_h + 0.1

    return s


def slide_activity(prs, title, steps, time_str=None):
    """Activity slide with numbered steps and teal header."""
    s = _blank(prs)
    _bg(s, S_CREAM)
    _header_bar(s, title, fill=S_TEAL, time_str=time_str)

    max_steps = min(len(steps), 5)
    step_h = min(0.55, 3.8 / max(max_steps, 1))
    y = 1.1
    for i, step in enumerate(steps[:max_steps], 1):
        _rect(s, Inches(0.6), Inches(y), Inches(0.40), Inches(0.40), S_TEAL)
        _txt(s, str(i), Inches(0.6), Inches(y + 0.02), Inches(0.40), Inches(0.36),
             font=SERIF, sz=16, bold=True, color=S_WHITE, align=PP_ALIGN.CENTER)
        _rich_txt(s, step[:150], Inches(1.15), Inches(y + 0.05), Inches(8.2), Inches(0.35),
                  font=SANS, sz=14, color=S_TEXT)
        y += step_h + 0.1

    return s


def slide_discussion(prs, title, questions):
    """Discussion slide with large question numbers."""
    s = _blank(prs)
    _bg(s, S_CREAM)
    _header_bar(s, title, fill=S_NAVY)

    n = len(questions[:3])
    row_h = min(1.2, 4.0 / max(n, 1))
    y = 1.15

    for i, q in enumerate(questions[:3], 1):
        _txt(s, str(i), Inches(0.5), Inches(y), Inches(0.5), Inches(0.45),
             font=SERIF, sz=24, bold=True, color=S_TEAL, align=PP_ALIGN.CENTER)
        _rich_txt(s, q[:200], Inches(1.1), Inches(y + 0.02), Inches(8.2), Inches(row_h - 0.15),
                  font=SANS, sz=16, color=S_TEXT)
        y += row_h

    return s


def slide_exit_ticket(prs, questions):
    """Exit ticket slide."""
    s = _blank(prs)
    _bg(s, S_WHITE)
    _header_bar(s, "Exit Ticket", fill=S_NAVY)
    _txt(s, "Complete independently before you leave.",
         Inches(0.6), Inches(1.0), Inches(8.8), Inches(0.3),
         font=SANS, sz=13, italic=True, color=S_GRAY)

    n = len(questions[:3])
    card_h = min(1.1, 3.5 / max(n, 1))
    y = 1.4

    for i, q in enumerate(questions[:3], 1):
        _rect(s, Inches(0.6), Inches(y), Inches(0.06), Inches(card_h), S_TEAL)
        _rect(s, Inches(0.66), Inches(y), Inches(8.74), Inches(card_h), S_LIGHT)
        _txt(s, f"Question {i}", Inches(0.9), Inches(y + 0.06), Inches(8.3), Inches(0.25),
             font=SANS, sz=11, bold=True, color=S_TEAL)
        _rich_txt(s, q[:250], Inches(0.9), Inches(y + 0.30), Inches(8.3), Inches(card_h - 0.38),
                  font=SANS, sz=14, color=S_TEXT)
        y += card_h + 0.15

    return s


def slide_divider(prs, title, subtitle=""):
    """Section divider (dark background)."""
    s = _blank(prs)
    _bg(s, S_DARK_BG)
    _txt(s, title, Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.0),
         font=SERIF, sz=30, bold=True, color=S_WHITE)
    if subtitle:
        _txt(s, subtitle, Inches(0.8), Inches(2.9), Inches(8.4), Inches(0.6),
             font=SANS, sz=16, color=S_GRAY)
    _rect(s, Inches(0.8), Inches(3.8), Inches(1.5), Inches(0.04), S_TEAL)
    return s


# ═══════════════════════════════════════════════════════════════════════════════
# CONTENT PARSING
# ═══════════════════════════════════════════════════════════════════════════════

def _parse_sections(content: str) -> Dict[str, str]:
    """Parse ## headers into {name: body} dict."""
    sections = {}
    lines = content.split("\n")
    key = None
    buf = []

    for line in lines:
        if line.startswith("## "):
            if key is not None:
                sections[key] = "\n".join(buf).strip()
            key = line[3:].strip()
            buf = []
        elif not line.startswith("# "):
            buf.append(line)

    if key is not None:
        sections[key] = "\n".join(buf).strip()

    return sections


def _get_bullets(text: str, n: int = 6) -> List[str]:
    """Extract bullet/numbered items from text."""
    items = []
    for line in text.split("\n"):
        line = line.strip()
        m = re.match(r'^[-*\u2022]\s+(.+)', line)
        if m:
            items.append(m.group(1).strip())
            continue
        m = re.match(r'^\d+[.)]\s+(.+)', line)
        if m:
            items.append(m.group(1).strip())
    return items[:n]


def _get_questions(text: str) -> List[str]:
    """Extract numbered questions."""
    qs = []
    cur = []
    for line in text.split("\n"):
        line = line.strip()
        m = re.match(r'^\d+[.)]\s+(.+)', line)
        if m:
            if cur:
                qs.append(" ".join(cur))
            cur = [m.group(1)]
        elif line.startswith("-") or line.startswith("*"):
            continue  # skip follow-ups for slide
        elif line and cur:
            cur.append(line)
    if cur:
        qs.append(" ".join(cur))
    return qs


def _strip_answer_lines(text: str) -> str:
    """Remove lines that look like answer keys from student-facing content."""
    cleaned = []
    for line in text.split("\n"):
        stripped = line.strip().lower()
        if any(stripped.startswith(p) for p in [
            "answer:", "correct answer:", "correct:", "solution:",
            "expected response:", "expected answer:", "sample answer:",
            "key:", "answer key:", "**answer", "**correct",
            "**expected", "**solution",
        ]):
            continue
        # Skip bare letter answers like "B" or "(B)"
        if re.match(r'^[\(\[]?[A-Da-d][\)\]]?\.?$', stripped):
            continue
        cleaned.append(line)
    return "\n".join(cleaned)


def _get_time(text: str) -> Optional[str]:
    """Extract time reference like '5 min'."""
    m = re.search(r'(\d+)\s*(?:min(?:ute)?s?)', text, re.IGNORECASE)
    return f"{m.group(1)} min" if m else None


def _get_bold_terms(text: str) -> List[Tuple[str, str]]:
    """Extract **term**: definition patterns."""
    terms = []
    for line in text.split("\n"):
        m = re.match(r'^[-*\u2022]?\s*\*\*(.+?)\*\*[:\u2014\u2013\-]+\s*(.+)', line)
        if m:
            terms.append((m.group(1).strip(), m.group(2).strip()))
    return terms


def _get_paras(text: str, n: int = 5) -> List[str]:
    """Get non-empty paragraphs, truncated for slide use."""
    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    result = []
    for p in paras[:n]:
        clean = p.replace("\n", " ").strip()
        if clean:
            result.append(clean[:150])
    return result


def _split_lessons(content: str) -> List[str]:
    """Split multi-lesson content on '# Lesson N' headers."""
    parts = re.split(r'(?=^#\s+Lesson\s+\d)', content, flags=re.MULTILINE)
    return [p.strip() for p in parts if p.strip()]


# ═══════════════════════════════════════════════════════════════════════════════
# MASTER SLIDE DECK BUILDER
# ═══════════════════════════════════════════════════════════════════════════════

def _build_lesson_slides(prs, sections: Dict[str, str]):
    """Build slides for one lesson's worth of content sections."""

    # Extract chart blocks from all sections upfront
    cleaned_sections = {}
    all_charts = {}
    for key, raw in sections.items():
        cleaned, charts = _extract_chart_blocks(raw)
        cleaned_sections[key] = cleaned
        if charts:
            all_charts[key] = charts

    # Do Now
    if "Do Now" in sections:
        raw = cleaned_sections.get("Do Now", sections["Do Now"])
        items = _get_bullets(raw) or _get_paras(raw, 3)
        slide_content(prs, "Do Now", items, notes_text=raw,
                      time_str=_get_time(raw) or "5 min", bg_color=S_WHITE)
        for chart in all_charts.get("Do Now", []):
            slide_visual(prs, chart)

    # Framing
    if "Framing" in sections:
        raw = cleaned_sections.get("Framing", sections["Framing"])
        items = _get_bullets(raw) or _get_paras(raw, 4)
        slide_content(prs, "Why This Matters", items,
                      notes_text=raw, time_str="3 min")
        for chart in all_charts.get("Framing", []):
            slide_visual(prs, chart)

    # Core Content — vocabulary cards + charts + tables + concept bullets
    if "Core Content" in sections:
        raw = cleaned_sections.get("Core Content", sections["Core Content"])
        vocab = _get_bold_terms(raw)
        if vocab:
            slide_cards(prs, "Key Vocabulary",
                        [(t, d) for t, d in vocab[:4]])

        # Render any chart visuals
        for chart in all_charts.get("Core Content", []):
            slide_visual(prs, chart)

        # Detect and render markdown tables as slide tables
        table_lines = [l for l in raw.split("\n") if l.strip().startswith("|")]
        if table_lines:
            headers, tbl_rows = _parse_md_table(table_lines)
            if headers and tbl_rows:
                slide_table(prs, "Key Concepts", headers, tbl_rows)

        # Non-table content as bullet slides
        non_table_lines = [l for l in raw.split("\n") if not l.strip().startswith("|")]
        non_table_raw = "\n".join(non_table_lines)
        items = _get_bullets(non_table_raw)
        if not items:
            items = _get_paras(non_table_raw, 5)
        # Split into slides of max 5 items
        for chunk_start in range(0, len(items), 5):
            chunk = items[chunk_start:chunk_start + 5]
            if chunk:
                lbl = "Key Concepts" if chunk_start == 0 else "Key Concepts (continued)"
                slide_content(prs, lbl, chunk, notes_text=sections["Core Content"])

    # Discussion
    if "Discussion Prompts" in sections:
        raw = sections["Discussion Prompts"]
        qs = _get_questions(raw)
        if qs:
            slide_discussion(prs, "Discussion", qs[:3])

    # Activity
    if "Activity Instructions" in sections:
        raw = sections["Activity Instructions"]
        steps = _get_bullets(raw)
        t = _get_time(raw)
        if steps:
            slide_activity(prs, "Activity", steps[:5], time_str=t)
        else:
            paras = _get_paras(raw, 4)
            slide_content(prs, "Activity", paras,
                          notes_text=raw, time_str=t, bar_color=S_TEAL)

    # Exit Ticket (strip any answer lines that leaked through)
    if "Exit Ticket" in sections:
        raw = _strip_answer_lines(sections["Exit Ticket"])
        qs = _get_questions(raw) or _get_bullets(raw)
        if not qs:
            qs = [raw[:250]]
        slide_exit_ticket(prs, qs[:3])


def build_slide_deck(content: str, title: str, objectives: List[str],
                     duration: int, output_path: str):
    """Build a complete professional slide deck from Phase 4 content."""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    # Check for multi-lesson format
    lessons = _split_lessons(content) if re.search(r'^#\s+Lesson\s+\d', content, re.MULTILINE) else None

    if lessons and len(lessons) > 1:
        # Multi-lesson: get plan from first lesson or full content
        all_sections = _parse_sections(content)
        plan = all_sections.get("Lesson Plan", content[:3000])

        slide_hidden_plan(prs, plan)
        slide_title(prs, title, f"{len(lessons)}-Lesson Sequence \u2022 {duration} min each")
        slide_objectives(prs, objectives, duration)

        for i, lesson_content in enumerate(lessons):
            obj = objectives[i] if i < len(objectives) else objectives[-1]
            slide_divider(prs, f"Lesson {i+1}", obj)
            sections = _parse_sections(lesson_content)
            _build_lesson_slides(prs, sections)
    else:
        # Single lesson
        sections = _parse_sections(content)
        plan = sections.get("Lesson Plan", content[:3000])

        slide_hidden_plan(prs, plan)
        slide_title(prs, title, f"{duration}-Minute Lesson")
        slide_objectives(prs, objectives, duration)
        _build_lesson_slides(prs, sections)

    prs.save(output_path)


# ═══════════════════════════════════════════════════════════════════════════════
# WORKSHEET HELPERS (python-docx)
# ═══════════════════════════════════════════════════════════════════════════════

def _cell_shading(cell, color):
    shd = OxmlElement('w:shd')
    shd.set(docx_qn('w:fill'), color)
    shd.set(docx_qn('w:val'), 'clear')
    cell._tc.get_or_add_tcPr().append(shd)


def _cell_borders(cell, color=D_BORDER, sz="1"):
    tc_pr = cell._tc.get_or_add_tcPr()
    borders = OxmlElement('w:tcBorders')
    for name in ['top', 'left', 'bottom', 'right']:
        b = OxmlElement(f'w:{name}')
        b.set(docx_qn('w:val'), 'single')
        b.set(docx_qn('w:sz'), sz)
        b.set(docx_qn('w:color'), color)
        borders.append(b)
    tc_pr.append(borders)


def _cell_margins(cell, top=60, bottom=60, left=100, right=100):
    tc_pr = cell._tc.get_or_add_tcPr()
    margins = OxmlElement('w:tcMar')
    for attr, val in [('top', top), ('left', left), ('bottom', bottom), ('right', right)]:
        m = OxmlElement(f'w:{attr}')
        m.set(docx_qn('w:w'), str(val))
        m.set(docx_qn('w:type'), 'dxa')
        margins.append(m)
    tc_pr.append(margins)


def _cell_valign(cell, val='center'):
    tc_pr = cell._tc.get_or_add_tcPr()
    va = OxmlElement('w:vAlign')
    va.set(docx_qn('w:val'), val)
    tc_pr.append(va)


def _doc_para(doc, text, font=SANS, sz=10, bold=False, italic=False,
              color=None, align=None, before=None, after=None):
    """Add a formatted paragraph to the document. Parses **bold** markdown."""
    p = doc.add_paragraph()
    text = str(text)

    # Parse **bold** markdown into runs
    parts = re.split(r'(\*\*.*?\*\*)', text)
    has_markdown = any(part.startswith("**") and part.endswith("**") for part in parts if part)

    if has_markdown:
        for part in parts:
            if not part:
                continue
            if part.startswith("**") and part.endswith("**"):
                run = p.add_run(part[2:-2])
                run.font.bold = True
            else:
                run = p.add_run(part)
                if bold:
                    run.font.bold = True
                if italic:
                    run.font.italic = True
            run.font.name = font
            run.font.size = DocxPt(sz)
            if color:
                run.font.color.rgb = _drgb(color)
    else:
        run = p.add_run(text)
        run.font.name = font
        run.font.size = DocxPt(sz)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = _drgb(color)

    if align:
        p.alignment = align
    pf = p.paragraph_format
    if before is not None:
        pf.space_before = DocxEmu(before)
    if after is not None:
        pf.space_after = DocxEmu(after)
    return p


def _cell_text(cell, text, font=SANS, sz=10, bold=False, color=D_TEXT, align=None):
    """Set formatted text in a table cell."""
    p = cell.paragraphs[0]
    p.clear()
    run = p.add_run(str(text))
    run.font.name = font
    run.font.size = DocxPt(sz)
    run.font.bold = bold
    if color:
        run.font.color.rgb = _drgb(color)
    if align:
        p.alignment = align


def _styled_table(doc, headers: List[str], rows: List[List[str]],
                  col_widths: Optional[List[int]] = None):
    """Add a table with colored header row and alternating row shading."""
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    table = doc.add_table(rows=n_rows, cols=n_cols)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    # Calculate column widths
    if col_widths is None:
        total = 10800  # 7.5 inches in twips
        col_widths = [total // n_cols] * n_cols

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        _cell_shading(cell, D_NAVY)
        _cell_borders(cell)
        _cell_margins(cell)
        _cell_valign(cell)
        _cell_text(cell, header, font=SANS, sz=10, bold=True,
                   color=D_WHITE, align=WD_ALIGN_PARAGRAPH.CENTER)
        cell.width = Twips(col_widths[j])

    # Data rows with alternating shading
    for i, row_data in enumerate(rows):
        shade = D_LAVENDER if i % 2 == 0 else D_WHITE
        for j in range(n_cols):
            cell = table.cell(i + 1, j)
            _cell_shading(cell, shade)
            _cell_borders(cell)
            _cell_margins(cell)
            text = row_data[j] if j < len(row_data) else ""
            is_first_col = (j == 0)
            _cell_text(cell, text, font=SANS, sz=10,
                       bold=is_first_col, color=D_TEXT)
            cell.width = Twips(col_widths[j])

    # Add small spacing after table
    doc.add_paragraph().paragraph_format.space_after = DocxEmu(25400)

    return table


def _response_table(doc, n_lines=3, label=""):
    """Add a bordered response area using a table (instead of underscores)."""
    table = doc.add_table(rows=n_lines, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    width = 10800

    for i in range(n_lines):
        cell = table.cell(i, 0)
        shade = D_LAVENDER if i % 2 == 0 else D_WHITE
        _cell_shading(cell, shade)
        _cell_borders(cell)
        _cell_margins(cell, top=80, bottom=80)
        _cell_text(cell, " ", font=SANS, sz=10)
        cell.width = Twips(width)

    doc.add_paragraph().paragraph_format.space_after = DocxEmu(12700)
    return table


# ═══════════════════════════════════════════════════════════════════════════════
# WORKSHEET CONTENT PARSING
# ═══════════════════════════════════════════════════════════════════════════════

def _parse_md_table(lines: List[str]) -> Tuple[List[str], List[List[str]]]:
    """Parse markdown table lines into (headers, rows)."""
    headers = []
    rows = []
    for line in lines:
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if all(re.match(r'^[-:]+$', c.strip()) for c in cells if c.strip()):
            continue
        if not headers:
            headers = cells
        else:
            rows.append(cells)
    return headers, rows


def _parse_worksheet_parts(text: str) -> List[Dict]:
    """Parse worksheet content into structured parts."""
    parts = []
    current = None

    for line in text.split("\n"):
        # Detect part headers: ### Part N:, ### Section, **Part N:**
        m = re.match(r'^###\s+(.+)', line)
        if not m:
            m = re.match(r'^\*\*\s*(Part\s+\d.+?)\*\*', line)
        if m:
            if current:
                parts.append(current)
            current = {"title": m.group(1).strip().rstrip(":"), "elements": []}
            continue

        if current is None:
            continue

        stripped = line.strip()

        # Table row
        if stripped.startswith("|") and "|" in stripped[1:]:
            if current["elements"] and current["elements"][-1]["type"] == "table":
                current["elements"][-1]["rows"].append(line)
            else:
                current["elements"].append({"type": "table", "rows": [line]})
            continue

        # Numbered question
        qm = re.match(r'^(\d+)[.)]\s+(.+)', stripped)
        if qm:
            current["elements"].append({
                "type": "question",
                "number": qm.group(1),
                "text": qm.group(2),
            })
            continue

        # Underscore response line — skip (we use table cells)
        if re.match(r'^[_\s]{5,}$', stripped):
            continue

        # Regular text
        if stripped:
            current["elements"].append({"type": "text", "content": stripped})

    if current:
        parts.append(current)

    return parts


# ═══════════════════════════════════════════════════════════════════════════════
# MASTER WORKSHEET BUILDER
# ═══════════════════════════════════════════════════════════════════════════════

def _setup_doc() -> Document:
    """Create a Document with US Letter page setup and 0.5" margins."""
    doc = Document()
    section = doc.sections[0]
    section.page_width = DocxInches(8.5)
    section.page_height = DocxInches(11)
    section.top_margin = DocxInches(0.5)
    section.bottom_margin = DocxInches(0.5)
    section.left_margin = DocxInches(0.5)
    section.right_margin = DocxInches(0.5)
    return doc


def _add_doc_header(doc, title, badge=None):
    """Add worksheet title and name/date line."""
    # Title
    _doc_para(doc, title, font=SERIF, sz=14, bold=True, color=D_NAVY,
              align=WD_ALIGN_PARAGRAPH.CENTER, after=25400)

    # Badge (for modified/extension)
    if badge:
        _doc_para(doc, badge, font=SANS, sz=9, bold=True, color=D_TEAL,
                  align=WD_ALIGN_PARAGRAPH.RIGHT, after=25400)

    # Name / Date / Period line
    _doc_para(doc, "Name: ________________________    Date: ____________    Period: ______",
              font=SANS, sz=10, color=D_GRAY,
              align=WD_ALIGN_PARAGRAPH.CENTER, after=127000)


def _build_part(doc, part: Dict, body_sz=10):
    """Build one worksheet part (section header + elements)."""
    # Section header
    is_first = (part == part)  # always true — spacing handled below
    _doc_para(doc, part["title"], font=SERIF, sz=12, bold=True, color=D_NAVY,
              before=190500, after=63500)

    for elem in part["elements"]:
        if elem["type"] == "text":
            # Detect instruction-like text (starts with verb or "For each...")
            text = elem["content"]
            is_instruction = (text[0].isupper() and
                              any(text.lower().startswith(w) for w in
                                  ["for ", "use ", "look ", "read ", "write ", "complete ",
                                   "answer ", "based ", "using ", "in ", "after ", "before ",
                                   "predict ", "explain ", "describe ", "compare ", "analyze ",
                                   "list ", "identify ", "think ", "reflect ", "consider ",
                                   "choose ", "select ", "match ", "fill ", "circle ",
                                   "underline ", "highlight "]))
            if is_instruction:
                _doc_para(doc, text, font=SANS, sz=body_sz, italic=True,
                          color=D_GRAY, after=76200)
            else:
                _doc_para(doc, text, font=SANS, sz=body_sz, after=38100)

        elif elem["type"] == "table":
            headers, rows = _parse_md_table(elem["rows"])
            if headers:
                # Calculate column widths based on content
                n = len(headers)
                total = 10800
                # Give first column less space if it's a number column
                if headers[0].strip() in ("#", "No.", "No", "#.", ""):
                    first_w = 700
                    rest_w = (total - first_w) // max(n - 1, 1)
                    widths = [first_w] + [rest_w] * (n - 1)
                else:
                    widths = [total // n] * n
                _styled_table(doc, headers, rows, col_widths=widths)

        elif elem["type"] == "question":
            # Question with response area
            q_text = f"{elem['number']}. {elem['text']}"
            _doc_para(doc, q_text, font=SANS, sz=body_sz, bold=False, after=38100)
            # Only add response table if question doesn't already have inline blanks
            has_blanks = ("____" in elem["text"] or "______" in elem["text"]
                          or "___:" in elem["text"])
            if not has_blanks:
                _response_table(doc, n_lines=2)


def build_worksheet(content: str, title: str, output_path: str):
    """Build a standard student worksheet from Phase 4 content."""
    doc = _setup_doc()
    _add_doc_header(doc, title)

    sections = _parse_sections(content)

    # Primary source: ## Worksheet Content
    ws_text = sections.get("Worksheet Content", "")

    if ws_text:
        parts = _parse_worksheet_parts(ws_text)
        if parts:
            for part in parts:
                _build_part(doc, part)
        else:
            # Worksheet Content exists but has no ### sub-parts — render as flat content
            _doc_para(doc, "Complete each section as directed.",
                      font=SANS, sz=10, italic=True, color=D_GRAY, after=76200)
            bullets = _get_bullets(ws_text)
            if bullets:
                for b in bullets:
                    _doc_para(doc, f"\u2022 {b}", font=SANS, sz=10, after=38100)
                    _response_table(doc, n_lines=2)
            else:
                paras = _get_paras(ws_text, 10)
                for p_text in paras:
                    _doc_para(doc, p_text, font=SANS, sz=10, after=38100)
                    _response_table(doc, n_lines=2)
    else:
        # Fallback: build from other sections
        _doc_para(doc, "Complete each section as directed.",
                  font=SANS, sz=10, italic=True, color=D_GRAY, after=76200)

        for key in ["Do Now", "Core Content", "Activity Instructions", "Exit Ticket"]:
            if key in sections:
                _doc_para(doc, key, font=SERIF, sz=12, bold=True,
                          color=D_NAVY, before=190500, after=63500)
                bullets = _get_bullets(sections[key])
                if bullets:
                    for b in bullets[:5]:
                        _doc_para(doc, f"\u2022 {b}", font=SANS, sz=10, after=38100)
                _response_table(doc, n_lines=3)

    # Always add Exit Ticket at end if not already in worksheet parts
    et_text = _strip_answer_lines(sections.get("Exit Ticket", ""))
    if et_text and "Exit Ticket" not in ws_text:
        _doc_para(doc, "Exit Ticket", font=SERIF, sz=12, bold=True,
                  color=D_NAVY, before=190500, after=63500)
        qs = _get_questions(et_text) or _get_bullets(et_text)
        if qs:
            for i, q in enumerate(qs[:3], 1):
                _doc_para(doc, f"{i}. {q}", font=SANS, sz=10, after=38100)
                _response_table(doc, n_lines=2)
        else:
            _doc_para(doc, et_text[:500], font=SANS, sz=10, after=38100)
            _response_table(doc, n_lines=3)

    doc.save(output_path)


def build_modified_worksheet(content: str, title: str, output_path: str):
    """Build a modified worksheet for struggling learners.
    Adds word bank, increases font size, keeps same structure.
    """
    doc = _setup_doc()
    _add_doc_header(doc, title, badge="MODIFIED")

    sections = _parse_sections(content)

    # Extract vocabulary for word bank
    core = sections.get("Core Content", "")
    vocab = _get_bold_terms(core)
    if vocab:
        _doc_para(doc, "Word Bank \u2014 Use These Words to Help You",
                  font=SERIF, sz=11, bold=True, color=D_TEAL,
                  before=63500, after=50800)
        # Word bank in a shaded box
        words = ", ".join(t for t, _ in vocab)
        p = doc.add_paragraph()
        run = p.add_run(words)
        run.font.name = SANS
        run.font.size = DocxPt(10)
        run.font.bold = True
        run.font.color.rgb = _drgb(D_NAVY)
        p.paragraph_format.space_after = DocxEmu(127000)
        # Add shading to the paragraph
        pPr = p._element.get_or_add_pPr()
        shd = OxmlElement('w:shd')
        shd.set(docx_qn('w:fill'), D_LAVENDER)
        shd.set(docx_qn('w:val'), 'clear')
        pPr.append(shd)

    # Build worksheet parts with larger font
    ws_text = sections.get("Worksheet Content", "")
    body_sz = 11  # Larger for accessibility

    if ws_text:
        parts = _parse_worksheet_parts(ws_text)
        for part in parts:
            _build_part(doc, part, body_sz=body_sz)
    else:
        for key in ["Do Now", "Core Content", "Activity Instructions", "Exit Ticket"]:
            if key in sections:
                _doc_para(doc, key, font=SERIF, sz=12, bold=True,
                          color=D_NAVY, before=190500, after=63500)
                _response_table(doc, n_lines=3)

    # Exit Ticket
    et_text = _strip_answer_lines(sections.get("Exit Ticket", ""))
    if et_text and "Exit Ticket" not in ws_text:
        _doc_para(doc, "Exit Ticket", font=SERIF, sz=12, bold=True,
                  color=D_NAVY, before=190500, after=63500)
        qs = _get_questions(et_text) or _get_bullets(et_text)
        for i, q in enumerate((qs or [et_text[:300]])[:3], 1):
            _doc_para(doc, f"{i}. {q}", font=SANS, sz=body_sz, after=38100)
            _response_table(doc, n_lines=3)

    doc.save(output_path)


def build_extension_worksheet(content: str, title: str, output_path: str):
    """Build an extension worksheet for advanced learners.
    Same structure as standard + deeper challenge section.
    """
    doc = _setup_doc()
    _add_doc_header(doc, title, badge="EXTENSION")

    sections = _parse_sections(content)
    ws_text = sections.get("Worksheet Content", "")

    if ws_text:
        parts = _parse_worksheet_parts(ws_text)
        for part in parts:
            _build_part(doc, part)

    # Add challenge section
    _doc_para(doc, "Challenge: Deep Analysis", font=SERIF, sz=12, bold=True,
              color=D_NAVY, before=190500, after=63500)
    _doc_para(doc, "Answer the following questions using evidence and reasoning.",
              font=SANS, sz=10, italic=True, color=D_GRAY, after=76200)

    # Pull from discussion prompts for deeper questions
    disc = sections.get("Discussion Prompts", "")
    deep_qs = _get_questions(disc)
    if deep_qs:
        for i, q in enumerate(deep_qs[:3], 1):
            _doc_para(doc, f"{i}. {q}", font=SANS, sz=10, after=38100)
            _response_table(doc, n_lines=4)
    else:
        _doc_para(doc, "1. What are the most important connections between the concepts in today's lesson? Explain with specific evidence.",
                  font=SANS, sz=10, after=38100)
        _response_table(doc, n_lines=4)
        _doc_para(doc, "2. If you could design a follow-up investigation, what question would you explore and why?",
                  font=SANS, sz=10, after=38100)
        _response_table(doc, n_lines=4)

    # Exit Ticket
    et_text = _strip_answer_lines(sections.get("Exit Ticket", ""))
    if et_text and "Exit Ticket" not in ws_text:
        _doc_para(doc, "Exit Ticket", font=SERIF, sz=12, bold=True,
                  color=D_NAVY, before=190500, after=63500)
        qs = _get_questions(et_text) or _get_bullets(et_text)
        for i, q in enumerate((qs or [et_text[:300]])[:3], 1):
            _doc_para(doc, f"{i}. {q}", font=SANS, sz=10, after=38100)
            _response_table(doc, n_lines=2)

    doc.save(output_path)


def build_supplementary(content: str, title: str, output_path: str):
    """Build supplementary materials (station cards, data sheets) from Group Materials section."""
    doc = _setup_doc()
    _add_doc_header(doc, f"{title} \u2014 Supplementary Materials")

    sections = _parse_sections(content)
    group_text = sections.get("Group Materials", "")

    if not group_text:
        _doc_para(doc, "No supplementary materials specified for this lesson.",
                  font=SANS, sz=10, italic=True, color=D_GRAY, after=76200)
        doc.save(output_path)
        return

    # Split group materials by ### subsections (station cards, etc.)
    cards = []
    current_title = None
    current_body = []

    for line in group_text.split("\n"):
        m = re.match(r'^###\s+(.+)', line)
        if m:
            if current_title:
                cards.append((current_title, "\n".join(current_body).strip()))
            current_title = m.group(1).strip()
            current_body = []
        else:
            current_body.append(line)

    if current_title:
        cards.append((current_title, "\n".join(current_body).strip()))

    if not cards:
        # No ### headers — treat entire section as one card
        cards = [("Materials", group_text)]

    for card_title, card_body in cards:
        # Card header with teal background
        _doc_para(doc, card_title, font=SERIF, sz=12, bold=True, color=D_NAVY,
                  before=190500, after=63500)

        # Card content
        for line in card_body.split("\n"):
            line = line.strip()
            if not line:
                continue

            # Table rows
            if line.startswith("|"):
                continue  # Tables handled below

            # Bullets
            bm = re.match(r'^[-*\u2022]\s+(.+)', line)
            if bm:
                _doc_para(doc, f"\u2022 {bm.group(1)}", font=SANS, sz=10, after=38100)
                continue

            # Numbered items
            nm = re.match(r'^\d+[.)]\s+(.+)', line)
            if nm:
                _doc_para(doc, line, font=SANS, sz=10, after=38100)
                _response_table(doc, n_lines=2)
                continue

            # Bold headers
            if line.startswith("**") and line.endswith("**"):
                _doc_para(doc, line[2:-2], font=SANS, sz=10, bold=True, after=38100)
                continue

            # Regular text
            _doc_para(doc, line, font=SANS, sz=10, after=38100)

        # Check for tables in card body
        table_lines = [l for l in card_body.split("\n") if l.strip().startswith("|")]
        if table_lines:
            headers, rows = _parse_md_table(table_lines)
            if headers:
                _styled_table(doc, headers, rows)

        # Page break between cards (except last)
        if card_title != cards[-1][0]:
            doc.add_page_break()

    doc.save(output_path)
